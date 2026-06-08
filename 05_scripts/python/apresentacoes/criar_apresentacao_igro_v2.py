"""
Gerador de Apresentação IGRO v2 — COM METODOLOGIA E RESULTADOS
10 slides máximo
Baseado em Guias de Estilo CGE-GO
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import datetime

# ========== CORES OFICIAIS CGE-GO ==========
CORES = {
    'verde_oficial': RGBColor(31, 162, 46),
    'azul_oficial': RGBColor(0, 81, 158),
    'teal': RGBColor(0, 118, 111),
    'verde_escuro': RGBColor(5, 66, 34),
    'laranja': RGBColor(247, 147, 30),
    'laranja_escuro': RGBColor(213, 111, 36),
    'amarelo_quente': RGBColor(251, 176, 59),
    'branco': RGBColor(255, 255, 255),
    'cinza_claro': RGBColor(245, 245, 245),
    'cinza_texto': RGBColor(74, 74, 74),
    'verde_sucesso': RGBColor(28, 173, 71),
}

# ========== FUNÇÕES AUXILIARES ==========

def adicionar_fundo_cor(slide, cor_rgb):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = cor_rgb

def adicionar_titulo(slide, texto, tamanho=44, cor=CORES['verde_escuro']):
    left = Inches(0.5)
    top = Inches(0.5)
    width = Inches(9)
    height = Inches(1)

    title_box = slide.shapes.add_textbox(left, top, width, height)
    title_frame = title_box.text_frame
    title_frame.word_wrap = True

    p = title_frame.paragraphs[0]
    p.text = texto
    p.font.size = Pt(tamanho)
    p.font.bold = True
    p.font.color.rgb = cor
    p.font.name = 'Segoe UI'

    return title_box

def adicionar_linha_separadora(slide, top, cor=CORES['teal'], width=10, left=0):
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(0.02)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = cor
    line.line.color.rgb = cor

def adicionar_kpi_card(slide, valor, label, left, top, cor_valor=CORES['teal']):
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(2), Inches(1.2)
    )
    card.fill.solid()
    card.fill.fore_color.rgb = CORES['branco']
    card.line.color.rgb = CORES['cinza_claro']
    card.line.width = Pt(1)

    valor_box = slide.shapes.add_textbox(Inches(left), Inches(top + 0.1), Inches(2), Inches(0.6))
    valor_frame = valor_box.text_frame
    p = valor_frame.paragraphs[0]
    p.text = str(valor)
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = cor_valor
    p.font.name = 'Segoe UI'
    p.alignment = PP_ALIGN.CENTER

    label_box = slide.shapes.add_textbox(Inches(left), Inches(top + 0.65), Inches(2), Inches(0.4))
    label_frame = label_box.text_frame
    p = label_frame.paragraphs[0]
    p.text = label
    p.font.size = Pt(10)
    p.font.color.rgb = CORES['cinza_texto']
    p.font.name = 'Segoe UI'
    p.alignment = PP_ALIGN.CENTER

    return card

def adicionar_footer(slide):
    footer_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(6.8), Inches(10), Inches(0.7)
    )
    footer_shape.fill.solid()
    footer_shape.fill.fore_color.rgb = CORES['cinza_claro']
    footer_shape.line.color.rgb = CORES['cinza_claro']

    footer_box = slide.shapes.add_textbox(Inches(0.3), Inches(6.85), Inches(9.4), Inches(0.6))
    footer_frame = footer_box.text_frame

    p = footer_frame.paragraphs[0]
    p.text = "IGRO — Controladoria-Geral do Estado de Goiás | Maio de 2026"
    p.font.size = Pt(9)
    p.font.color.rgb = CORES['cinza_texto']
    p.font.name = 'Segoe UI'

# ========== CRIAR APRESENTAÇÃO ==========

prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# ========== SLIDE 1: CAPA ==========
slide1 = prs.slides.add_slide(prs.slide_layouts[6])
adicionar_fundo_cor(slide1, CORES['verde_escuro'])

stripe = slide1.shapes.add_shape(
    MSO_SHAPE.RECTANGLE,
    Inches(9), Inches(0), Inches(1), Inches(7.5)
)
stripe.fill.solid()
stripe.fill.fore_color.rgb = RGBColor(255, 221, 0)
stripe.line.color.rgb = RGBColor(255, 221, 0)

brasao = slide1.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(2), Inches(0.8))
brasao_frame = brasao.text_frame
p = brasao_frame.paragraphs[0]
p.text = "CGE-GO"
p.font.size = Pt(28)
p.font.color.rgb = CORES['branco']
p.font.bold = True

titulo1 = slide1.shapes.add_textbox(Inches(0.5), Inches(2), Inches(8), Inches(2))
titulo_frame = titulo1.text_frame
titulo_frame.word_wrap = True
p = titulo_frame.paragraphs[0]
p.text = "Índice de Gestão de Riscos de Ouvidoria"
p.font.size = Pt(54)
p.font.bold = True
p.font.color.rgb = CORES['branco']
p.font.name = 'Segoe UI'
p.line_spacing = 1.2

subtitulo1 = slide1.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(8), Inches(1.5))
subtitulo_frame = subtitulo1.text_frame
subtitulo_frame.word_wrap = True
p = subtitulo_frame.paragraphs[0]
p.text = "IGRO — Metodologia para Monitoramento Estratégico"
p.font.size = Pt(20)
p.font.color.rgb = CORES['amarelo_quente']
p.font.name = 'Segoe UI'

info = slide1.shapes.add_textbox(Inches(0.5), Inches(6.2), Inches(8), Inches(0.8))
info_frame = info.text_frame
p = info_frame.paragraphs[0]
p.text = "Maio de 2026 | Alta Gestão"
p.font.size = Pt(12)
p.font.color.rgb = CORES['branco']
p.font.name = 'Segoe UI'

# ========== SLIDE 2: O PROBLEMA ==========
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
adicionar_fundo_cor(slide2, CORES['branco'])

stripe2 = slide2.shapes.add_shape(
    MSO_SHAPE.RECTANGLE,
    Inches(0), Inches(0), Inches(0.1), Inches(7.5)
)
stripe2.fill.solid()
stripe2.fill.fore_color.rgb = CORES['verde_escuro']
stripe2.line.color.rgb = CORES['verde_escuro']

adicionar_titulo(slide2, "O PROBLEMA", tamanho=40)
adicionar_linha_separadora(slide2, 1.1, cor=CORES['teal'])

corpo = slide2.shapes.add_textbox(Inches(0.7), Inches(1.5), Inches(8.6), Inches(5))
corpo_frame = corpo.text_frame
corpo_frame.word_wrap = True

linhas = [
    "109.338 manifestações cidadãs (2024-2025)",
    "",
    "Dados operacionais FRAGMENTADOS em múltiplos indicadores",
    "",
    "Sinais CONFLITANTES entre dimensões",
    "• Órgão rápido ≠ órgão efetivo",
    "• Falta visão estratégica unificada",
    "",
    "HETEROGENEIDADE extrema entre 51 órgãos",
]

for i, linha in enumerate(linhas):
    if i == 0:
        p = corpo_frame.paragraphs[0]
    else:
        p = corpo_frame.add_paragraph()

    p.text = linha
    p.font.size = Pt(13)
    p.font.color.rgb = CORES['cinza_texto']
    p.font.name = 'Segoe UI'
    p.level = 1 if linha.startswith("•") else 0
    p.space_before = Pt(4)

    if "109" in linha or "FRAGMENTADOS" in linha or "CONFLITANTES" in linha or "HETEROGENEIDADE" in linha:
        p.font.bold = True
        p.font.color.rgb = CORES['verde_escuro']

adicionar_footer(slide2)

# ========== SLIDE 3: METODOLOGIA ==========
slide3 = prs.slides.add_slide(prs.slide_layouts[6])
adicionar_fundo_cor(slide3, CORES['branco'])

stripe3 = slide3.shapes.add_shape(
    MSO_SHAPE.RECTANGLE,
    Inches(0), Inches(0), Inches(0.1), Inches(7.5)
)
stripe3.fill.solid()
stripe3.fill.fore_color.rgb = CORES['verde_escuro']
stripe3.line.color.rgb = CORES['verde_escuro']

adicionar_titulo(slide3, "METODOLOGIA", tamanho=40)
adicionar_linha_separadora(slide3, 1.1, cor=CORES['teal'])

# Referencial
ref = slide3.shapes.add_textbox(Inches(0.7), Inches(1.4), Inches(8.6), Inches(0.35))
rf = ref.text_frame
p = rf.paragraphs[0]
p.text = "Referencial: OCDE/JRC Handbook | ISO 31000:2018 | COSO"
p.font.size = Pt(11)
p.font.color.rgb = CORES['teal']
p.font.bold = True

# EIXO 1
eixo1_box = slide3.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(0.7), Inches(2), Inches(4), Inches(2.5)
)
eixo1_box.fill.solid()
eixo1_box.fill.fore_color.rgb = CORES['cinza_claro']
eixo1_box.line.color.rgb = CORES['teal']
eixo1_box.line.width = Pt(2)

t1 = slide3.shapes.add_textbox(Inches(0.9), Inches(2.1), Inches(3.6), Inches(0.3))
t1f = t1.text_frame
p = t1f.paragraphs[0]
p.text = "EIXO 1: TEMPESTIVIDADE"
p.font.size = Pt(12)
p.font.bold = True
p.font.color.rgb = CORES['teal']

c1 = slide3.shapes.add_textbox(Inches(0.9), Inches(2.5), Inches(3.6), Inches(1.8))
c1f = c1.text_frame
c1f.word_wrap = True
items1 = [
    "TMR: Tempo Médio de Resposta",
    "Meta: <= 5 dias (excelência)",
    "",
    "PMA: Manifestações em Atraso",
    "Meta: <= 1% (conformidade)"
]
for i, item in enumerate(items1):
    if i == 0:
        p = c1f.paragraphs[0]
    else:
        p = c1f.add_paragraph()
    p.text = item
    p.font.size = Pt(10)
    p.font.color.rgb = CORES['cinza_texto']
    p.space_before = Pt(2)

# EIXO 2
eixo2_box = slide3.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(5.3), Inches(2), Inches(4), Inches(2.5)
)
eixo2_box.fill.solid()
eixo2_box.fill.fore_color.rgb = CORES['cinza_claro']
eixo2_box.line.color.rgb = CORES['verde_oficial']
eixo2_box.line.width = Pt(2)

t2 = slide3.shapes.add_textbox(Inches(5.5), Inches(2.1), Inches(3.6), Inches(0.3))
t2f = t2.text_frame
p = t2f.paragraphs[0]
p.text = "EIXO 2: QUALIDADE"
p.font.size = Pt(12)
p.font.bold = True
p.font.color.rgb = CORES['verde_oficial']

c2 = slide3.shapes.add_textbox(Inches(5.5), Inches(2.5), Inches(3.6), Inches(1.8))
c2f = c2.text_frame
c2f.word_wrap = True
items2 = [
    "RP: Resolutividade Percebida",
    "NR: Nota de Recomendação",
    "%RI: Respostas Insatisfatórias",
    "",
    "Mede: Efetividade percebida"
]
for i, item in enumerate(items2):
    if i == 0:
        p = c2f.paragraphs[0]
    else:
        p = c2f.add_paragraph()
    p.text = item
    p.font.size = Pt(10)
    p.font.color.rgb = CORES['cinza_texto']
    p.space_before = Pt(2)

# Agregação
agr = slide3.shapes.add_textbox(Inches(0.7), Inches(4.7), Inches(8.6), Inches(1.8))
agf = agr.text_frame
agf.word_wrap = True

p = agf.paragraphs[0]
p.text = "AGREGACAO: Media Geometrica Ponderada"
p.font.size = Pt(12)
p.font.bold = True
p.font.color.rgb = CORES['verde_escuro']

p = agf.add_paragraph()
p.text = "Formula: IGRO = (KRI1 x KRI2 x KRI3 x KRI4 x KRI5)^(1/5)"
p.font.size = Pt(10)
p.font.color.rgb = CORES['cinza_texto']
p.space_before = Pt(6)

p = agf.add_paragraph()
p.text = "Beneficio: Penaliza fragilidades extremas; força equilíbrio"
p.font.size = Pt(10)
p.font.color.rgb = CORES['cinza_texto']
p.space_before = Pt(4)

adicionar_footer(slide3)

# ========== SLIDE 4: RESULTADOS CONSOLIDADOS ==========
slide4 = prs.slides.add_slide(prs.slide_layouts[6])
adicionar_fundo_cor(slide4, CORES['branco'])

stripe4 = slide4.shapes.add_shape(
    MSO_SHAPE.RECTANGLE,
    Inches(0), Inches(0), Inches(0.1), Inches(7.5)
)
stripe4.fill.solid()
stripe4.fill.fore_color.rgb = CORES['verde_escuro']
stripe4.line.color.rgb = CORES['verde_escuro']

adicionar_titulo(slide4, "RESULTADOS CONSOLIDADOS", tamanho=38)
adicionar_linha_separadora(slide4, 1.1, cor=CORES['teal'])

# KPI Cards
adicionar_kpi_card(slide4, "52,9%", "IGRO\n(Crítico)", 0.5, 1.6, CORES['laranja_escuro'])
adicionar_kpi_card(slide4, "6,8 dias", "TMR", 2.8, 1.6, CORES['teal'])
adicionar_kpi_card(slide4, "2,3%", "PMA\n(Atraso)", 5.1, 1.6, CORES['laranja'])
adicionar_kpi_card(slide4, "61,5%", "Resolutiv.", 7.4, 1.6, CORES['verde_oficial'])

# Interpretação
interp = slide4.shapes.add_textbox(Inches(0.7), Inches(3.2), Inches(8.6), Inches(3.2))
iaf = interp.text_frame
iaf.word_wrap = True

p = iaf.paragraphs[0]
p.text = "Diagnostico:"
p.font.size = Pt(12)
p.font.bold = True
p.font.color.rgb = CORES['verde_escuro']

items_interp = [
    "IGRO 52,9% = NIVEL CRITICO (média rede opera com risco elevado)",
    "TMR 6,8 dias = Abaixo prazo legal (bom), mas HETEROGENEO entre órgãos",
    "PMA 2,3% = Não conformidade parcial (alguns órgãos descumprem prazo)",
    "Resolutividade 61,5% = Moderada (4 em 10 cidadãos com dúvida/insatisfação)"
]

for item in items_interp:
    p = iaf.add_paragraph()
    p.text = "• " + item
    p.font.size = Pt(10)
    p.font.color.rgb = CORES['cinza_texto']
    p.space_before = Pt(3)

adicionar_footer(slide4)

# ========== SLIDE 5: ANALISE DE RESULTADOS ==========
slide5 = prs.slides.add_slide(prs.slide_layouts[6])
adicionar_fundo_cor(slide5, CORES['branco'])

stripe5 = slide5.shapes.add_shape(
    MSO_SHAPE.RECTANGLE,
    Inches(0), Inches(0), Inches(0.1), Inches(7.5)
)
stripe5.fill.solid()
stripe5.fill.fore_color.rgb = CORES['verde_escuro']
stripe5.line.color.rgb = CORES['verde_escuro']

adicionar_titulo(slide5, "ANALISE DE RESULTADOS", tamanho=38)
adicionar_linha_separadora(slide5, 1.1, cor=CORES['teal'])

analise = slide5.shapes.add_textbox(Inches(0.7), Inches(1.4), Inches(8.6), Inches(5.1))
anf = analise.text_frame
anf.word_wrap = True

items_analise = [
    ("HETEROGENEIDADE EXTREMA", "GOINFRA (100%) vs SECAMI (~30%) — não é variação normal, é diferença estrutural"),
    ("ORGAOS CRITICOS = MAIS VOLUME", "39% dos órgãos em nível crítico; justamente ali há MAIS manifestações cidadãs"),
    ("TAMANHO NAO EXPLICA DESEMPENHO", "Órgãos pequenos não são sempre ruins; órgãos grandes não são sempre bons"),
    ("DISSOCIACAO TMR vs RESOLUTIVIDADE", "Alguns órgãos rápidos têm baixa resolutividade; outros lentos mas efetivos"),
]

for i, (titulo, desc) in enumerate(items_analise):
    if i == 0:
        p = anf.paragraphs[0]
    else:
        p = anf.add_paragraph()

    p.text = titulo
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = CORES['teal']
    p.space_before = Pt(6)

    p = anf.add_paragraph()
    p.text = "→ " + desc
    p.font.size = Pt(10)
    p.font.color.rgb = CORES['cinza_texto']
    p.level = 1

adicionar_footer(slide5)

# ========== SLIDE 6: LEI DE GOODHART ==========
slide6 = prs.slides.add_slide(prs.slide_layouts[6])
adicionar_fundo_cor(slide6, CORES['branco'])

stripe6 = slide6.shapes.add_shape(
    MSO_SHAPE.RECTANGLE,
    Inches(0), Inches(0), Inches(0.1), Inches(7.5)
)
stripe6.fill.solid()
stripe6.fill.fore_color.rgb = CORES['verde_escuro']
stripe6.line.color.rgb = CORES['verde_escuro']

adicionar_titulo(slide6, "ACHADO CRITICO", tamanho=40)
adicionar_titulo(slide6, "Lei de Goodhart", tamanho=24, cor=CORES['laranja'])

citacao_box = slide6.shapes.add_textbox(Inches(0.7), Inches(1.8), Inches(8.6), Inches(1.2))
cit_frame = citacao_box.text_frame
cit_frame.word_wrap = True
p = cit_frame.paragraphs[0]
p.text = '"Quando uma medida se torna uma meta, ela deixa de ser uma boa medida."'
p.font.size = Pt(16)
p.font.italic = True
p.font.color.rgb = CORES['laranja']
p.font.bold = True
p.alignment = PP_ALIGN.CENTER

evidencia = slide6.shapes.add_textbox(Inches(0.7), Inches(3.1), Inches(8.6), Inches(0.5))
ev_frame = evidencia.text_frame
p = ev_frame.paragraphs[0]
p.text = "Observacao nos dados: Orgaos com TMR baixo nem sempre tiveram Resolutividade alta"
p.font.size = Pt(12)
p.font.color.rgb = CORES['verde_escuro']
p.alignment = PP_ALIGN.CENTER

implicacao = slide6.shapes.add_textbox(Inches(0.7), Inches(3.8), Inches(8.6), Inches(2.7))
imp_frame = implicacao.text_frame
imp_frame.word_wrap = True

linhas_imp = [
    "RISCO: Otimizacao exclusiva de velocidade",
    "→ Orgao responde rapido mas superficial",
    "",
    "SOLUCAO: Indicador composto (IGRO)",
    "→ Força equilibrio entre velocidade E qualidade",
    "→ Impossivel compensar fragilidade em uma dimensao",
    "→ Mitiga comportamentos disfuncionais"
]

for i, linha in enumerate(linhas_imp):
    if i == 0:
        p = imp_frame.paragraphs[0]
    else:
        p = imp_frame.add_paragraph()

    p.text = linha
    p.font.size = Pt(11)
    p.font.color.rgb = CORES['cinza_texto']
    p.level = 1 if "→" in linha else 0
    p.space_before = Pt(3)

adicionar_footer(slide6)

# ========== SLIDE 7: FATORES DE SUCESSO ==========
slide7 = prs.slides.add_slide(prs.slide_layouts[6])
adicionar_fundo_cor(slide7, CORES['branco'])

stripe7 = slide7.shapes.add_shape(
    MSO_SHAPE.RECTANGLE,
    Inches(0), Inches(0), Inches(0.1), Inches(7.5)
)
stripe7.fill.solid()
stripe7.fill.fore_color.rgb = CORES['verde_escuro']
stripe7.line.color.rgb = CORES['verde_escuro']

adicionar_titulo(slide7, "5 FATORES DE SUCESSO", tamanho=38)
adicionar_linha_separadora(slide7, 1.1, cor=CORES['teal'])

fatores = [
    ("Integracao Tecnologica", "Sistemas modernos = resposta rapida + estabilidade"),
    ("Capacidade Tecnica Dedicada", "Equipes especializadas = consistencia nos KPIs"),
    ("Comunicacao Estruturada", "Linguagem clara = percepcao cidada positiva"),
    ("Conformidade Normativa", "Processos padronizados = reduz surpresas"),
]

y_pos = 1.5
for titulo, desc in fatores:
    fator_box = slide7.shapes.add_textbox(Inches(0.7), Inches(y_pos), Inches(8.6), Inches(0.25))
    ff = fator_box.text_frame
    p = ff.paragraphs[0]
    p.text = f"• {titulo}"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = CORES['teal']

    desc_box = slide7.shapes.add_textbox(Inches(1.1), Inches(y_pos + 0.28), Inches(8.2), Inches(0.35))
    df = desc_box.text_frame
    df.word_wrap = True
    p = df.paragraphs[0]
    p.text = desc
    p.font.size = Pt(10)
    p.font.color.rgb = CORES['cinza_texto']

    y_pos += 0.88

msg = slide7.shapes.add_textbox(Inches(0.7), Inches(5.8), Inches(8.6), Inches(0.7))
mf = msg.text_frame
mf.word_wrap = True
p = mf.paragraphs[0]
p.text = "REPLICAVEL: Sao escolhas de gestao, nao destino. Orgaos criticos podem adotar esses 5 fatores."
p.font.size = Pt(11)
p.font.bold = True
p.font.color.rgb = CORES['verde_escuro']

adicionar_footer(slide7)

# ========== SLIDE 8: APLICACOES PRATICAS ==========
slide8 = prs.slides.add_slide(prs.slide_layouts[6])
adicionar_fundo_cor(slide8, CORES['branco'])

stripe8 = slide8.shapes.add_shape(
    MSO_SHAPE.RECTANGLE,
    Inches(0), Inches(0), Inches(0.1), Inches(7.5)
)
stripe8.fill.solid()
stripe8.fill.fore_color.rgb = CORES['verde_escuro']
stripe8.line.color.rgb = CORES['verde_escuro']

adicionar_titulo(slide8, "APLICACOES PRATICAS", tamanho=38)
adicionar_linha_separadora(slide8, 1.1, cor=CORES['teal'])

aplicacoes = [
    ("Monitoramento Estrategico", "Dashboard executivo, IGRO atualizado mensal/trimestral"),
    ("Avaliacao Comparativa", "Benchmarking entre orgaos, replicacao de boas praticas"),
    ("Priorizacao de Acoes", "Recursos direcionados para orgaos criticos"),
    ("Governanca por Risco", "Integracao com Matriz de Riscos da CGE"),
]

y_pos = 1.5
for titulo, desc in aplicacoes:
    app_box = slide8.shapes.add_textbox(Inches(0.7), Inches(y_pos), Inches(8.6), Inches(0.25))
    af = app_box.text_frame
    p = af.paragraphs[0]
    p.text = f"• {titulo}"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = CORES['verde_oficial']

    desc_box = slide8.shapes.add_textbox(Inches(1.1), Inches(y_pos + 0.28), Inches(8.2), Inches(0.35))
    df = desc_box.text_frame
    df.word_wrap = True
    p = df.paragraphs[0]
    p.text = desc
    p.font.size = Pt(10)
    p.font.color.rgb = CORES['cinza_texto']

    y_pos += 0.88

adicionar_footer(slide8)

# ========== SLIDE 9: PROXIMOS PASSOS ==========
slide9 = prs.slides.add_slide(prs.slide_layouts[6])
adicionar_fundo_cor(slide9, CORES['branco'])

stripe9 = slide9.shapes.add_shape(
    MSO_SHAPE.RECTANGLE,
    Inches(0), Inches(0), Inches(0.1), Inches(7.5)
)
stripe9.fill.solid()
stripe9.fill.fore_color.rgb = CORES['verde_escuro']
stripe9.line.color.rgb = CORES['verde_escuro']

adicionar_titulo(slide9, "PROXIMOS PASSOS", tamanho=38)
adicionar_linha_separadora(slide9, 1.1, cor=CORES['teal'])

# Fase 1
fase1_box = slide9.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(0.4), Inches(1.5), Inches(3), Inches(4.8)
)
fase1_box.fill.solid()
fase1_box.fill.fore_color.rgb = CORES['cinza_claro']
fase1_box.line.color.rgb = CORES['laranja']
fase1_box.line.width = Pt(2)

f1_titulo = slide9.shapes.add_textbox(Inches(0.6), Inches(1.65), Inches(2.6), Inches(0.3))
f1tf = f1_titulo.text_frame
p = f1tf.paragraphs[0]
p.text = "FASE 1 — Imediato"
p.font.size = Pt(11)
p.font.bold = True
p.font.color.rgb = CORES['laranja']

f1_items = [
    "Validar metodologia",
    "Comunicar rede",
    "Estruturar dashboard"
]
y_f1 = 2.05
for item in f1_items:
    f1_item = slide9.shapes.add_textbox(Inches(0.7), Inches(y_f1), Inches(2.4), Inches(0.3))
    f1if = f1_item.text_frame
    p = f1if.paragraphs[0]
    p.text = "✓ " + item
    p.font.size = Pt(10)
    p.font.color.rgb = CORES['cinza_texto']
    y_f1 += 0.6

# Fase 2
fase2_box = slide9.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(3.5), Inches(1.5), Inches(3), Inches(4.8)
)
fase2_box.fill.solid()
fase2_box.fill.fore_color.rgb = CORES['cinza_claro']
fase2_box.line.color.rgb = CORES['amarelo_quente']
fase2_box.line.width = Pt(2)

f2_titulo = slide9.shapes.add_textbox(Inches(3.7), Inches(1.65), Inches(2.6), Inches(0.3))
f2tf = f2_titulo.text_frame
p = f2tf.paragraphs[0]
p.text = "FASE 2 — 3-6 meses"
p.font.size = Pt(11)
p.font.bold = True
p.font.color.rgb = CORES['amarelo_quente']

f2_items = [
    "Implementar calculos",
    "Capacitar gestores",
    "Planos de melhoria"
]
y_f2 = 2.05
for item in f2_items:
    f2_item = slide9.shapes.add_textbox(Inches(3.7), Inches(y_f2), Inches(2.6), Inches(0.3))
    f2if = f2_item.text_frame
    p = f2if.paragraphs[0]
    p.text = "✓ " + item
    p.font.size = Pt(10)
    p.font.color.rgb = CORES['cinza_texto']
    y_f2 += 0.6

# Fase 3
fase3_box = slide9.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(6.6), Inches(1.5), Inches(3), Inches(4.8)
)
fase3_box.fill.solid()
fase3_box.fill.fore_color.rgb = CORES['cinza_claro']
fase3_box.line.color.rgb = CORES['verde_oficial']
fase3_box.line.width = Pt(2)

f3_titulo = slide9.shapes.add_textbox(Inches(6.8), Inches(1.65), Inches(2.6), Inches(0.3))
f3tf = f3_titulo.text_frame
p = f3tf.paragraphs[0]
p.text = "FASE 3 — 6-12 meses"
p.font.size = Pt(11)
p.font.bold = True
p.font.color.rgb = CORES['verde_oficial']

f3_items = [
    "Novos indicadores",
    "Modelos preditivos",
    "Validacao longitudinal"
]
y_f3 = 2.05
for item in f3_items:
    f3_item = slide9.shapes.add_textbox(Inches(6.8), Inches(y_f3), Inches(2.6), Inches(0.3))
    f3if = f3_item.text_frame
    p = f3if.paragraphs[0]
    p.text = "✓ " + item
    p.font.size = Pt(10)
    p.font.color.rgb = CORES['cinza_texto']
    y_f3 += 0.6

adicionar_footer(slide9)

# ========== SLIDE 10: CONCLUSAO ==========
slide10 = prs.slides.add_slide(prs.slide_layouts[6])
adicionar_fundo_cor(slide10, CORES['verde_escuro'])

stripe10 = slide10.shapes.add_shape(
    MSO_SHAPE.RECTANGLE,
    Inches(9), Inches(0), Inches(1), Inches(7.5)
)
stripe10.fill.solid()
stripe10.fill.fore_color.rgb = RGBColor(255, 221, 0)
stripe10.line.color.rgb = RGBColor(255, 221, 0)

titulo_conclusao = slide10.shapes.add_textbox(Inches(0.5), Inches(1), Inches(8.5), Inches(1))
tc_frame = titulo_conclusao.text_frame
tc_frame.word_wrap = True
p = tc_frame.paragraphs[0]
p.text = "IGRO: De Dados a Decisao Estrategica"
p.font.size = Pt(40)
p.font.bold = True
p.font.color.rgb = CORES['branco']

checkmarks_texto = slide10.shapes.add_textbox(Inches(0.5), Inches(2.3), Inches(8.5), Inches(3.5))
ck_frame = checkmarks_texto.text_frame
ck_frame.word_wrap = True

checks = [
    "✓ Sintetiza 5 metricas em 1 numero estrategico",
    "✓ Força equilibrio entre velocidade e qualidade",
    "✓ Identifica heterogeneidade e oportunidades",
    "✓ Reduz comportamentos disfuncionais",
    "✓ Pronto para adocao imediata",
]

for i, check in enumerate(checks):
    if i == 0:
        p = ck_frame.paragraphs[0]
    else:
        p = ck_frame.add_paragraph()

    p.text = check
    p.font.size = Pt(13)
    p.font.color.rgb = CORES['branco']
    p.space_before = Pt(8)

# Chamada para acao
acao = slide10.shapes.add_textbox(Inches(0.5), Inches(6.2), Inches(8.5), Inches(0.8))
acao_frame = acao.text_frame
acao_frame.word_wrap = True
p = acao_frame.paragraphs[0]
p.text = "Vocês aprovam a validação e implementação da FASE 1?"
p.font.size = Pt(16)
p.font.bold = True
p.font.color.rgb = CORES['amarelo_quente']
p.alignment = PP_ALIGN.CENTER

# ========== SALVAR ==========
output_path = r"C:\Users\andrei.lima\OneDrive\Claude-Work\Projects\igro\APRESENTACAO_IGRO_v2.pptx"
prs.save(output_path)

print("Apresentacao criada: APRESENTACAO_IGRO_v2.pptx (10 slides)")
print("Com: Metodologia + Resultados detalhados")
