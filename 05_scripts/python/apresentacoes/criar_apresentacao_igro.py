"""
Gerador de Apresentação IGRO
Baseado em Guias de Estilo CGE-GO

Cores Oficiais:
- VERDE_OFICIAL: #1fa22e
- AZUL_OFICIAL: #00519e
- TEAL: #00766f (primária para dados)
- VERDE_ESCURO: #054222 (fundo premium)
- LARANJA: #f7931e (alertas/metas)
- CINZA_TEXTO: #4a4a4a
- BRANCO: #ffffff
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import datetime

# ========== CORES OFICIAIS CGE-GO ==========
CORES = {
    'verde_oficial': RGBColor(31, 162, 46),      # #1fa22e
    'azul_oficial': RGBColor(0, 81, 158),        # #00519e
    'teal': RGBColor(0, 118, 111),               # #00766f
    'verde_escuro': RGBColor(5, 66, 34),         # #054222
    'laranja': RGBColor(247, 147, 30),           # #f7931e
    'laranja_escuro': RGBColor(213, 111, 36),    # #d56f24
    'amarelo_quente': RGBColor(251, 176, 59),    # #fbb03b
    'branco': RGBColor(255, 255, 255),
    'cinza_claro': RGBColor(245, 245, 245),
    'cinza_texto': RGBColor(74, 74, 74),
    'verde_sucesso': RGBColor(28, 173, 71),      # #1CAD47
}

# ========== FUNÇÕES AUXILIARES ==========

def adicionar_fundo_cor(slide, cor_rgb):
    """Adiciona cor de fundo sólido ao slide"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = cor_rgb

def adicionar_titulo(slide, texto, tamanho=44, cor=CORES['verde_escuro']):
    """Adiciona título formatado"""
    left = Inches(0.5)
    top = Inches(0.5)
    width = Inches(9)
    height = Inches(1)

    title_box = slide.shapes.add_textbox(left, top, width, height)
    title_frame = title_box.text_frame
    title_frame.word_wrap = True

    p = title_frame.paragraphs[0]
    p.text = texto
    p.font.size = Pt(tamanho)
    p.font.bold = True
    p.font.color.rgb = cor
    p.font.name = 'Segoe UI'

    return title_box

def adicionar_subtitulo(slide, texto, left=0.5, top=1.2, tamanho=20, cor=CORES['teal']):
    """Adiciona subtítulo formatado"""
    width = Inches(9)
    height = Inches(0.6)

    subtitle_box = slide.shapes.add_textbox(Inches(left), Inches(top), width, height)
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.word_wrap = True

    p = subtitle_frame.paragraphs[0]
    p.text = texto
    p.font.size = Pt(tamanho)
    p.font.color.rgb = cor
    p.font.name = 'Segoe UI'

    return subtitle_box

def adicionar_corpo(slide, texto, left=0.5, top=2.0, width=9, height=4.5, tamanho=14, cor=CORES['cinza_texto']):
    """Adiciona texto de corpo"""
    text_box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    text_frame = text_box.text_frame
    text_frame.word_wrap = True

    p = text_frame.paragraphs[0]
    p.text = texto
    p.font.size = Pt(tamanho)
    p.font.color.rgb = cor
    p.font.name = 'Segoe UI'
    p.line_spacing = 1.4

    return text_box

def adicionar_kpi_card(slide, valor, label, left, top, cor_valor=CORES['teal']):
    """Adiciona card de KPI"""
    # Background do card
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(2), Inches(1.2)
    )
    card.fill.solid()
    card.fill.fore_color.rgb = CORES['branco']
    card.line.color.rgb = CORES['cinza_claro']
    card.line.width = Pt(1)

    # Valor grande
    valor_box = slide.shapes.add_textbox(Inches(left), Inches(top + 0.1), Inches(2), Inches(0.6))
    valor_frame = valor_box.text_frame
    p = valor_frame.paragraphs[0]
    p.text = str(valor)
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = cor_valor
    p.font.name = 'Segoe UI'
    p.alignment = PP_ALIGN.CENTER

    # Label
    label_box = slide.shapes.add_textbox(Inches(left), Inches(top + 0.65), Inches(2), Inches(0.4))
    label_frame = label_box.text_frame
    p = label_frame.paragraphs[0]
    p.text = label
    p.font.size = Pt(10)
    p.font.color.rgb = CORES['cinza_texto']
    p.font.name = 'Segoe UI'
    p.alignment = PP_ALIGN.CENTER

    return card

def adicionar_footer(slide):
    """Adiciona rodapé institucional"""
    footer_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(6.8), Inches(10), Inches(0.7)
    )
    footer_shape.fill.solid()
    footer_shape.fill.fore_color.rgb = CORES['cinza_claro']
    footer_shape.line.color.rgb = CORES['cinza_claro']

    # Texto footer
    footer_box = slide.shapes.add_textbox(Inches(0.3), Inches(6.85), Inches(9.4), Inches(0.6))
    footer_frame = footer_box.text_frame

    p = footer_frame.paragraphs[0]
    p.text = f"Índice de Gestão de Riscos de Ouvidoria (IGRO) | Controladoria-Geral do Estado de Goiás | Maio de 2026"
    p.font.size = Pt(9)
    p.font.color.rgb = CORES['cinza_texto']
    p.font.name = 'Segoe UI'

def adicionar_linha_separadora(slide, top, cor=CORES['teal'], width=10, left=0):
    """Adiciona linha separadora"""
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(0.02)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = cor
    line.line.color.rgb = cor

# ========== CRIAR APRESENTAÇÃO ==========

prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# ========== SLIDE 1: CAPA ==========
slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
adicionar_fundo_cor(slide1, CORES['verde_escuro'])

# Barra amarela (stripe)
stripe = slide1.shapes.add_shape(
    MSO_SHAPE.RECTANGLE,
    Inches(9), Inches(0), Inches(1), Inches(7.5)
)
stripe.fill.solid()
stripe.fill.fore_color.rgb = RGBColor(255, 221, 0)  # Amarelo oficial
stripe.line.color.rgb = RGBColor(255, 221, 0)

# Brasão (texto placeholder)
brasao = slide1.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(2), Inches(0.8))
brasao_frame = brasao.text_frame
p = brasao_frame.paragraphs[0]
p.text = "🏛️ CGE-GO"
p.font.size = Pt(36)
p.font.color.rgb = CORES['branco']

# Título
titulo1 = slide1.shapes.add_textbox(Inches(0.5), Inches(2), Inches(8), Inches(2))
titulo_frame = titulo1.text_frame
titulo_frame.word_wrap = True
p = titulo_frame.paragraphs[0]
p.text = "Índice de Gestão de Riscos de Ouvidoria"
p.font.size = Pt(54)
p.font.bold = True
p.font.color.rgb = CORES['branco']
p.font.name = 'Segoe UI'
p.line_spacing = 1.2

# Subtítulo
subtitulo1 = slide1.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(8), Inches(1.5))
subtitulo_frame = subtitulo1.text_frame
subtitulo_frame.word_wrap = True
p = subtitulo_frame.paragraphs[0]
p.text = "IGRO — Metodologia para Monitoramento Estratégico da Rede Estadual de Ouvidorias"
p.font.size = Pt(20)
p.font.color.rgb = CORES['amarelo_quente']
p.font.name = 'Segoe UI'

# Data/Info
info = slide1.shapes.add_textbox(Inches(0.5), Inches(6.2), Inches(8), Inches(0.8))
info_frame = info.text_frame
p = info_frame.paragraphs[0]
p.text = "Maio de 2026 | Apresentação para Alta Gestão"
p.font.size = Pt(12)
p.font.color.rgb = CORES['branco']
p.font.name = 'Segoe UI'

# ========== SLIDE 2: O PROBLEMA ==========
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
adicionar_fundo_cor(slide2, CORES['branco'])

# Stripe lateral
stripe2 = slide2.shapes.add_shape(
    MSO_SHAPE.RECTANGLE,
    Inches(0), Inches(0), Inches(0.1), Inches(7.5)
)
stripe2.fill.solid()
stripe2.fill.fore_color.rgb = CORES['verde_escuro']
stripe2.line.color.rgb = CORES['verde_escuro']

adicionar_titulo(slide2, "O PROBLEMA", tamanho=40)
adicionar_linha_separadora(slide2, 1.1, cor=CORES['teal'])

# Conteúdo
corpo = slide2.shapes.add_textbox(Inches(0.7), Inches(1.5), Inches(8.6), Inches(5))
corpo_frame = corpo.text_frame
corpo_frame.word_wrap = True

linhas = [
    "109.338 manifestações recebidas entre 2024 e 2025",
    "",
    "Dados operacionais fragmentados em múltiplos indicadores",
    "",
    "❌ Sinais conflitantes (tempo rápido ≠ qualidade)",
    "❌ Falta de visão estratégica unificada",
    "❌ Impossível comparar órgãos com estruturas diferentes",
    "",
    "Pergunta central: Como transformar dados em decisão?",
]

for i, linha in enumerate(linhas):
    if i == 0:
        p = corpo_frame.paragraphs[0]
    else:
        p = corpo_frame.add_paragraph()

    p.text = linha
    p.font.size = Pt(14 if linha.startswith("❌") or linha.startswith("109") else 13)
    p.font.color.rgb = CORES['cinza_texto']
    p.font.name = 'Segoe UI'
    p.level = 1 if linha.startswith("❌") else 0
    p.space_before = Pt(6)

    if linha.startswith("Pergunta"):
        p.font.bold = True
        p.font.color.rgb = CORES['verde_escuro']

adicionar_footer(slide2)

# ========== SLIDE 3: A SOLUÇÃO (IGRO) ==========
slide3 = prs.slides.add_slide(prs.slide_layouts[6])
adicionar_fundo_cor(slide3, CORES['branco'])

stripe3 = slide3.shapes.add_shape(
    MSO_SHAPE.RECTANGLE,
    Inches(0), Inches(0), Inches(0.1), Inches(7.5)
)
stripe3.fill.solid()
stripe3.fill.fore_color.rgb = CORES['verde_escuro']
stripe3.line.color.rgb = CORES['verde_escuro']

adicionar_titulo(slide3, "SOLUÇÃO: O IGRO", tamanho=40)
adicionar_linha_separadora(slide3, 1.1, cor=CORES['teal'])

# Explicação
explicacao = slide3.shapes.add_textbox(Inches(0.7), Inches(1.5), Inches(8.6), Inches(1.2))
expl_frame = explicacao.text_frame
expl_frame.word_wrap = True
p = expl_frame.paragraphs[0]
p.text = "Um único indicador que sintetiza 5 métricas em 2 eixos, forçando equilíbrio entre velocidade e qualidade"
p.font.size = Pt(14)
p.font.color.rgb = CORES['verde_escuro']
p.font.bold = True
p.font.name = 'Segoe UI'

# Caixa TEMPESTIVIDADE
caixa_temp = slide3.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(0.7), Inches(2.9), Inches(4.2), Inches(3)
)
caixa_temp.fill.solid()
caixa_temp.fill.fore_color.rgb = CORES['cinza_claro']
caixa_temp.line.color.rgb = CORES['teal']
caixa_temp.line.width = Pt(2)

titulo_temp = slide3.shapes.add_textbox(Inches(0.9), Inches(3), Inches(3.8), Inches(0.4))
tf_temp = titulo_temp.text_frame
p = tf_temp.paragraphs[0]
p.text = "⚡ TEMPESTIVIDADE"
p.font.size = Pt(13)
p.font.bold = True
p.font.color.rgb = CORES['teal']

conteudo_temp = slide3.shapes.add_textbox(Inches(0.9), Inches(3.5), Inches(3.8), Inches(2))
cf_temp = conteudo_temp.text_frame
cf_temp.word_wrap = True
items_temp = [
    "TMR: Tempo Médio de Resposta",
    "PMA: Manifestações em Atraso"
]
for i, item in enumerate(items_temp):
    if i == 0:
        p = cf_temp.paragraphs[0]
    else:
        p = cf_temp.add_paragraph()
    p.text = item
    p.font.size = Pt(11)
    p.font.color.rgb = CORES['cinza_texto']
    p.space_before = Pt(4)

# Caixa QUALIDADE
caixa_qual = slide3.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(5.1), Inches(2.9), Inches(4.2), Inches(3)
)
caixa_qual.fill.solid()
caixa_qual.fill.fore_color.rgb = CORES['cinza_claro']
caixa_qual.line.color.rgb = CORES['verde_oficial']
caixa_qual.line.width = Pt(2)

titulo_qual = slide3.shapes.add_textbox(Inches(5.3), Inches(3), Inches(3.8), Inches(0.4))
tf_qual = titulo_qual.text_frame
p = tf_qual.paragraphs[0]
p.text = "✅ QUALIDADE"
p.font.size = Pt(13)
p.font.bold = True
p.font.color.rgb = CORES['verde_oficial']

conteudo_qual = slide3.shapes.add_textbox(Inches(5.3), Inches(3.5), Inches(3.8), Inches(2.2))
cf_qual = conteudo_qual.text_frame
cf_qual.word_wrap = True
items_qual = [
    "RP: Resolutividade Percebida",
    "%RI: Respostas Insatisfatórias",
    "NR: Nota de Recomendação (NPS)"
]
for i, item in enumerate(items_qual):
    if i == 0:
        p = cf_qual.paragraphs[0]
    else:
        p = cf_qual.add_paragraph()
    p.text = item
    p.font.size = Pt(11)
    p.font.color.rgb = CORES['cinza_texto']
    p.space_before = Pt(4)

adicionar_footer(slide3)

# ========== SLIDE 4: RESULTADOS PRINCIPAIS ==========
slide4 = prs.slides.add_slide(prs.slide_layouts[6])
adicionar_fundo_cor(slide4, CORES['branco'])

stripe4 = slide4.shapes.add_shape(
    MSO_SHAPE.RECTANGLE,
    Inches(0), Inches(0), Inches(0.1), Inches(7.5)
)
stripe4.fill.solid()
stripe4.fill.fore_color.rgb = CORES['verde_escuro']
stripe4.line.color.rgb = CORES['verde_escuro']

adicionar_titulo(slide4, "RESULTADOS PRINCIPAIS", tamanho=38)
adicionar_linha_separadora(slide4, 1.1, cor=CORES['teal'])

# KPI Cards
adicionar_kpi_card(slide4, "52,9%", "IGRO Consolidado\n(Crítico)", 0.5, 1.6, CORES['laranja_escuro'])
adicionar_kpi_card(slide4, "6,8 dias", "Tempo Médio\nde Resposta", 2.8, 1.6, CORES['teal'])
adicionar_kpi_card(slide4, "61,5%", "Resolutividade\nPercebida", 5.1, 1.6, CORES['verde_oficial'])
adicionar_kpi_card(slide4, "+32,4", "NPS\nModerado", 7.4, 1.6, CORES['amarelo_quente'])

# Achados
achados_titulo = slide4.shapes.add_textbox(Inches(0.7), Inches(3.3), Inches(8.6), Inches(0.4))
at_frame = achados_titulo.text_frame
p = at_frame.paragraphs[0]
p.text = "Achados-Chave:"
p.font.size = Pt(14)
p.font.bold = True
p.font.color.rgb = CORES['verde_escuro']

achados = [
    "39% dos órgãos em nível crítico de risco",
    "Forte heterogeneidade entre unidades (GOINFRA: 100% vs. SECAMI: ~30%)",
    "Volume operacional ≠ desempenho (fatores de gestão superam tamanho)"
]

y_pos = 3.8
for achado in achados:
    achado_box = slide4.shapes.add_textbox(Inches(1), Inches(y_pos), Inches(8.3), Inches(0.4))
    af = achado_box.text_frame
    p = af.paragraphs[0]
    p.text = "• " + achado
    p.font.size = Pt(12)
    p.font.color.rgb = CORES['cinza_texto']
    y_pos += 0.5

adicionar_footer(slide4)

# ========== SLIDE 5: LEI DE GOODHART (ACHADO CRÍTICO) ==========
slide5 = prs.slides.add_slide(prs.slide_layouts[6])
adicionar_fundo_cor(slide5, CORES['branco'])

stripe5 = slide5.shapes.add_shape(
    MSO_SHAPE.RECTANGLE,
    Inches(0), Inches(0), Inches(0.1), Inches(7.5)
)
stripe5.fill.solid()
stripe5.fill.fore_color.rgb = CORES['verde_escuro']
stripe5.line.color.rgb = CORES['verde_escuro']

adicionar_titulo(slide5, "ACHADO CRÍTICO", tamanho=38)
adicionar_titulo(slide5, "Lei de Goodhart", tamanho=24, cor=CORES['laranja'])

# Citação
citacao_box = slide5.shapes.add_textbox(Inches(0.7), Inches(1.8), Inches(8.6), Inches(1.2))
cit_frame = citacao_box.text_frame
cit_frame.word_wrap = True
p = cit_frame.paragraphs[0]
p.text = '"Quando uma medida se torna uma meta, ela deixa de ser uma boa medida."'
p.font.size = Pt(18)
p.font.italic = True
p.font.color.rgb = CORES['laranja']
p.font.bold = True
p.alignment = PP_ALIGN.CENTER

# Evidência
evidencia = slide5.shapes.add_textbox(Inches(0.7), Inches(3.2), Inches(8.6), Inches(0.5))
ev_frame = evidencia.text_frame
p = ev_frame.paragraphs[0]
p.text = "Observação: Órgãos com baixo TMR nem sempre tiveram alta Resolutividade Percebida"
p.font.size = Pt(13)
p.font.color.rgb = CORES['verde_escuro']
p.alignment = PP_ALIGN.CENTER

# Implicação
implicacao = slide5.shapes.add_textbox(Inches(0.7), Inches(4), Inches(8.6), Inches(2.5))
imp_frame = implicacao.text_frame
imp_frame.word_wrap = True

linhas_imp = [
    "⚠️  Otimização exclusiva de velocidade",
    "     → Respostas rápidas, porém superficiais",
    "",
    "✅ Solução: Indicador composto",
    "     → Força equilíbrio entre velocidade e qualidade",
    "     → Mitiga comportamentos disfuncionais"
]

for i, linha in enumerate(linhas_imp):
    if i == 0:
        p = imp_frame.paragraphs[0]
    else:
        p = imp_frame.add_paragraph()

    p.text = linha
    p.font.size = Pt(12)
    p.font.color.rgb = CORES['cinza_texto']
    p.level = 1 if "→" in linha else 0
    p.space_before = Pt(4)

adicionar_footer(slide5)

# ========== SLIDE 6: FATORES DE SUCESSO ==========
slide6 = prs.slides.add_slide(prs.slide_layouts[6])
adicionar_fundo_cor(slide6, CORES['branco'])

stripe6 = slide6.shapes.add_shape(
    MSO_SHAPE.RECTANGLE,
    Inches(0), Inches(0), Inches(0.1), Inches(7.5)
)
stripe6.fill.solid()
stripe6.fill.fore_color.rgb = CORES['verde_escuro']
stripe6.line.color.rgb = CORES['verde_escuro']

adicionar_titulo(slide6, "5 FATORES DE SUCESSO", tamanho=38)
adicionar_linha_separadora(slide6, 1.1, cor=CORES['teal'])

fatores = [
    ("💻", "Integração Tecnológica", "Sistemas modernos reduzem tempo e aumentam estabilidade"),
    ("👥", "Capacidade Técnica", "Equipes especializadas e treinamento contínuo"),
    ("📢", "Comunicação Estruturada", "Linguagem clara, acessível, centrada no cidadão"),
    ("✓", "Conformidade Normativa", "Processos padronizados e auditáveis"),
    ("📊", "Monitoramento Contínuo", "Acompanhamento regular e ajustes operacionais")
]

y_pos = 1.5
for icon, titulo, desc in fatores:
    # Icon + título
    fator_box = slide6.shapes.add_textbox(Inches(0.7), Inches(y_pos), Inches(8.6), Inches(0.3))
    ff = fator_box.text_frame
    p = ff.paragraphs[0]
    p.text = f"{icon}  {titulo}"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = CORES['teal']

    # Descrição
    desc_box = slide6.shapes.add_textbox(Inches(1.2), Inches(y_pos + 0.3), Inches(8.1), Inches(0.35))
    df = desc_box.text_frame
    df.word_wrap = True
    p = df.paragraphs[0]
    p.text = desc
    p.font.size = Pt(11)
    p.font.color.rgb = CORES['cinza_texto']

    y_pos += 1

adicionar_footer(slide6)

# ========== SLIDE 7: APLICAÇÕES PRÁTICAS ==========
slide7 = prs.slides.add_slide(prs.slide_layouts[6])
adicionar_fundo_cor(slide7, CORES['branco'])

stripe7 = slide7.shapes.add_shape(
    MSO_SHAPE.RECTANGLE,
    Inches(0), Inches(0), Inches(0.1), Inches(7.5)
)
stripe7.fill.solid()
stripe7.fill.fore_color.rgb = CORES['verde_escuro']
stripe7.line.color.rgb = CORES['verde_escuro']

adicionar_titulo(slide7, "APLICAÇÕES PRÁTICAS", tamanho=38)
adicionar_linha_separadora(slide7, 1.1, cor=CORES['teal'])

aplicacoes = [
    ("📊", "Monitoramento Estratégico", "Dashboard executivo com atualização mensal/trimestral"),
    ("🏆", "Avaliação Comparativa", "Benchmarking entre órgãos e identificação de boas práticas"),
    ("📍", "Priorização de Ações", "Direcionamento de recursos para órgãos críticos"),
    ("⚖️", "Governança orientada por Risco", "Integração com Matriz de Gestão de Riscos da CGE")
]

y_pos = 1.5
for icon, titulo, desc in aplicacoes:
    app_box = slide7.shapes.add_textbox(Inches(0.7), Inches(y_pos), Inches(8.6), Inches(0.3))
    af = app_box.text_frame
    p = af.paragraphs[0]
    p.text = f"{icon}  {titulo}"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = CORES['verde_oficial']

    desc_box = slide7.shapes.add_textbox(Inches(1.2), Inches(y_pos + 0.3), Inches(8.1), Inches(0.3))
    df = desc_box.text_frame
    df.word_wrap = True
    p = df.paragraphs[0]
    p.text = desc
    p.font.size = Pt(11)
    p.font.color.rgb = CORES['cinza_texto']

    y_pos += 1.1

adicionar_footer(slide7)

# ========== SLIDE 8: PRÓXIMOS PASSOS ==========
slide8 = prs.slides.add_slide(prs.slide_layouts[6])
adicionar_fundo_cor(slide8, CORES['branco'])

stripe8 = slide8.shapes.add_shape(
    MSO_SHAPE.RECTANGLE,
    Inches(0), Inches(0), Inches(0.1), Inches(7.5)
)
stripe8.fill.solid()
stripe8.fill.fore_color.rgb = CORES['verde_escuro']
stripe8.line.color.rgb = CORES['verde_escuro']

adicionar_titulo(slide8, "PRÓXIMOS PASSOS", tamanho=38)
adicionar_linha_separadora(slide8, 1.1, cor=CORES['teal'])

# Fase 1
fase1_box = slide8.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(0.5), Inches(1.5), Inches(2.8), Inches(4.5)
)
fase1_box.fill.solid()
fase1_box.fill.fore_color.rgb = CORES['cinza_claro']
fase1_box.line.color.rgb = CORES['laranja']
fase1_box.line.width = Pt(2)

f1_titulo = slide8.shapes.add_textbox(Inches(0.7), Inches(1.7), Inches(2.4), Inches(0.3))
f1tf = f1_titulo.text_frame
p = f1tf.paragraphs[0]
p.text = "FASE 1 — Imediato"
p.font.size = Pt(12)
p.font.bold = True
p.font.color.rgb = CORES['laranja']

f1_items = ["✓ Validar", "✓ Comunicar", "✓ Dashboards"]
y_f1 = 2.1
for item in f1_items:
    f1_item = slide8.shapes.add_textbox(Inches(0.9), Inches(y_f1), Inches(2.2), Inches(0.35))
    f1if = f1_item.text_frame
    p = f1if.paragraphs[0]
    p.text = item
    p.font.size = Pt(11)
    p.font.color.rgb = CORES['cinza_texto']
    y_f1 += 0.5

# Fase 2
fase2_box = slide8.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(3.6), Inches(1.5), Inches(2.8), Inches(4.5)
)
fase2_box.fill.solid()
fase2_box.fill.fore_color.rgb = CORES['cinza_claro']
fase2_box.line.color.rgb = CORES['amarelo_quente']
fase2_box.line.width = Pt(2)

f2_titulo = slide8.shapes.add_textbox(Inches(3.8), Inches(1.7), Inches(2.4), Inches(0.3))
f2tf = f2_titulo.text_frame
p = f2tf.paragraphs[0]
p.text = "FASE 2 — 3-6 meses"
p.font.size = Pt(12)
p.font.bold = True
p.font.color.rgb = CORES['amarelo_quente']

f2_items = ["✓ Implementar", "✓ Capacitar", "✓ Melhorar"]
y_f2 = 2.1
for item in f2_items:
    f2_item = slide8.shapes.add_textbox(Inches(3.8), Inches(y_f2), Inches(2.4), Inches(0.35))
    f2if = f2_item.text_frame
    p = f2if.paragraphs[0]
    p.text = item
    p.font.size = Pt(11)
    p.font.color.rgb = CORES['cinza_texto']
    y_f2 += 0.5

# Fase 3
fase3_box = slide8.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(6.7), Inches(1.5), Inches(2.8), Inches(4.5)
)
fase3_box.fill.solid()
fase3_box.fill.fore_color.rgb = CORES['cinza_claro']
fase3_box.line.color.rgb = CORES['verde_oficial']
fase3_box.line.width = Pt(2)

f3_titulo = slide8.shapes.add_textbox(Inches(6.9), Inches(1.7), Inches(2.4), Inches(0.3))
f3tf = f3_titulo.text_frame
p = f3tf.paragraphs[0]
p.text = "FASE 3 — 6-12 meses"
p.font.size = Pt(12)
p.font.bold = True
p.font.color.rgb = CORES['verde_oficial']

f3_items = ["✓ Evoluir", "✓ Integrar IA", "✓ Validar"]
y_f3 = 2.1
for item in f3_items:
    f3_item = slide8.shapes.add_textbox(Inches(6.9), Inches(y_f3), Inches(2.4), Inches(0.35))
    f3if = f3_item.text_frame
    p = f3if.paragraphs[0]
    p.text = item
    p.font.size = Pt(11)
    p.font.color.rgb = CORES['cinza_texto']
    y_f3 += 0.5

adicionar_footer(slide8)

# ========== SLIDE 9: CONCLUSÃO ==========
slide9 = prs.slides.add_slide(prs.slide_layouts[6])
adicionar_fundo_cor(slide9, CORES['verde_escuro'])

# Barra amarela
stripe9 = slide9.shapes.add_shape(
    MSO_SHAPE.RECTANGLE,
    Inches(9), Inches(0), Inches(1), Inches(7.5)
)
stripe9.fill.solid()
stripe9.fill.fore_color.rgb = RGBColor(255, 221, 0)
stripe9.line.color.rgb = RGBColor(255, 221, 0)

titulo_conclusao = slide9.shapes.add_textbox(Inches(0.5), Inches(1), Inches(8.5), Inches(0.8))
tc_frame = titulo_conclusao.text_frame
tc_frame.word_wrap = True
p = tc_frame.paragraphs[0]
p.text = "O IGRO oferece instrumento inédito"
p.font.size = Pt(36)
p.font.bold = True
p.font.color.rgb = CORES['branco']

subtitulo_conclusao = slide9.shapes.add_textbox(Inches(0.5), Inches(2), Inches(8.5), Inches(1))
sc_frame = subtitulo_conclusao.text_frame
sc_frame.word_wrap = True
p = sc_frame.paragraphs[0]
p.text = "para transformar dados operacionais de ouvidoria em linguagem estratégica"
p.font.size = Pt(20)
p.font.color.rgb = CORES['amarelo_quente']

# Checkmarks
checkmarks_texto = slide9.shapes.add_textbox(Inches(0.5), Inches(3.3), Inches(8.5), Inches(3))
ck_frame = checkmarks_texto.text_frame
ck_frame.word_wrap = True

checks = [
    "✓ Sintetiza múltiplas dimensões sem perder informação",
    "✓ Identifica heterogeneidade entre órgãos",
    "✓ Demonstra que gestão supera tamanho",
    "✓ Reduz comportamentos disfuncionais",
]

for i, check in enumerate(checks):
    if i == 0:
        p = ck_frame.paragraphs[0]
    else:
        p = ck_frame.add_paragraph()

    p.text = check
    p.font.size = Pt(14)
    p.font.color.rgb = CORES['branco']
    p.space_before = Pt(8)

# Chamada para ação
acao = slide9.shapes.add_textbox(Inches(0.5), Inches(6.2), Inches(8.5), Inches(0.8))
acao_frame = acao.text_frame
acao_frame.word_wrap = True
p = acao_frame.paragraphs[0]
p.text = "Pronto para adoção como ferramenta estratégica da CGE-GO"
p.font.size = Pt(16)
p.font.bold = True
p.font.color.rgb = CORES['verde_sucesso']
p.alignment = PP_ALIGN.CENTER

# ========== SALVAR ==========
output_path = r"C:\Users\andrei.lima\OneDrive\Claude-Work\Projects\igro\APRESENTACAO_IGRO_CGE-GO.pptx"
prs.save(output_path)

print(f"✅ Apresentação criada com sucesso!")
print(f"📁 Arquivo: {output_path}")
print(f"📊 Slides: 9")
print(f"🎨 Cores: Paleta oficial CGE-GO")
print(f"📝 Tipografia: Segoe UI")
